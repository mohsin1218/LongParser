"""Docling-based document extractor with Tesseract CLI OCR and HierarchicalChunker.

Uses:
- Tesseract CLI for OCR
- Layout analysis always enabled
- TableFormer for table structure
- HierarchicalChunker for heading hierarchy
- iterate_items() for reading-order block extraction

No hardcoded heuristics — relies entirely on Docling's native capabilities.
"""

from pathlib import Path
from typing import Optional, Tuple, List, Dict
import os
import time
import logging
import hashlib
import uuid
import re
from dataclasses import dataclass
from docling.datamodel.pipeline_options import (
    PdfPipelineOptions,
    TesseractCliOcrOptions,
)
from docling.datamodel.base_models import InputFormat
from docling.document_converter import (
    DocumentConverter,
    PdfFormatOption,
    WordFormatOption,
    PowerpointFormatOption,
    ExcelFormatOption,
    CsvFormatOption,
)

from docling_core.transforms.chunker import HierarchicalChunker
from docling_core.types.doc import (
    SectionHeaderItem,
    TableItem,
    PictureItem,
    TextItem,
    ListItem,
    DocItemLabel,
)

# TitleItem is used by Docling for PPTX slide titles (not SectionHeaderItem)
try:
    from docling_core.types.doc import TitleItem
except ImportError:
    TitleItem = None  # Fallback for older docling versions

from ..schemas import (
    Document, Page, Block, Table, TableCell,
    BlockType, ExtractorType, ProcessingConfig,
    BoundingBox, Provenance, Confidence, BlockFlags,
    DocumentMetadata, PageProfile, ExtractionMetadata,
)
from .base import BaseExtractor

logger = logging.getLogger(__name__)

# Pattern to detect structured leading markers in headings.
# Matches alphanumeric + punctuation prefixes followed by whitespace:
# "I.", "II.", "A.", "1.", "2.3", "IV", "a)", etc.
_MARKER_RE = re.compile(r'^([A-Za-z0-9][A-Za-z0-9.()]*)[.\s]\s*')

# Pattern used to detect garbled math in paragraph blocks.
_MATH_RE = re.compile(
    r'[\u2211\u220F\u222B\u221A\u00B1\u2264\u2265\u2248\u2260\u03B1-\u03C9\u03A3]'
    r'|[a-z]\s*=\s*[a-z0-9]',
    re.IGNORECASE,
)


def _iou_px(a: dict, b: dict) -> float:
    """Compute IoU between two pixel-space bbox dicts {x0,y0,x1,y1}."""
    xi0, yi0 = max(a["x0"], b["x0"]), max(a["y0"], b["y0"])
    xi1, yi1 = min(a["x1"], b["x1"]), min(a["y1"], b["y1"])
    inter = max(0, xi1 - xi0) * max(0, yi1 - yi0)
    ua = (a["x1"] - a["x0"]) * (a["y1"] - a["y0"])
    ub = (b["x1"] - b["x0"]) * (b["y1"] - b["y0"])
    union = ua + ub - inter
    return inter / union if union > 0 else 0.0


def _is_mfd_candidate(page_no: int, page_blocks, docling_formula_count: int) -> bool:
    """Return True if MFD should scan this page.

    Runs MFD if Docling found few/no formulas OR at least one non-equation
    block on this page contains garbled math Unicode.
    """
    if docling_formula_count > 3:
        return False  # Docling handled it well; trust it
    garbled = any(
        _MATH_RE.search(b.text)
        for b in page_blocks
        if getattr(b, "type", None) is not None and str(b.type) != "equation"
    )
    return docling_formula_count == 0 or garbled


@dataclass
class _HeadingInfo:
    """Internal heading tracking."""
    text: str
    level: int
    hierarchy_path: List[str]


@dataclass
class PptxParaInfo:
    """Paragraph info extracted directly from python-pptx."""
    indent_level: int      # 0-8 from paragraph.level
    is_title: bool         # True for TITLE / CENTER_TITLE placeholders
    is_subtitle: bool      # True for SUBTITLE placeholders
    is_list: bool          # True if Docling would treat it as list item
    bullet_type: str       # 'Bullet', 'Numbered', 'None'
    is_footer: bool = False  # True for DATE / FOOTER / SLIDE_NUMBER placeholders


@dataclass
class HierarchyChunk:
    """A chunk with hierarchy information."""
    text: str
    heading_path: List[str]
    level: int
    page_number: int
    order_index: int


class DoclingExtractor(BaseExtractor):
    """
    Document extractor using Docling with Tesseract CLI OCR.
    
    Relies entirely on Docling's native APIs:
    - iterate_items() for reading-order traversal with hierarchy level
    - SectionHeaderItem / TextItem / TableItem / ListItem / PictureItem for type detection
    - item.label (DocItemLabel) for fine-grained classification
    - item.prov for page number and bounding box
    - page.size for actual page dimensions
    - HierarchicalChunker for heading hierarchy paths
    
    Heading hierarchy is inferred autonomously from:
    1. Pattern Priority (Numbered vs Unnumbered)
    2. Position Awareness (Late Arrival Rule)
    3. Font-size clustering
    
    No hardcoded numbering conventions.
    """
    
    extractor_type = ExtractorType.DOCLING
    version = "3.0.0"
    
    def __init__(self, tesseract_lang: List[str] = None, tessdata_path: str = None, force_full_page_ocr: bool = False):
        """
        Initialize Docling extractor.
        
        Args:
            tesseract_lang: Languages for Tesseract OCR (default: ["eng"])
            tessdata_path: Path to tessdata directory with language models and configs.
                           If None, uses system default.
            force_full_page_ocr: If True, OCR entire page even if embedded text exists.
                                 Required for PDFs with broken Unicode mapping.
        """
        self._converter = None
        self._chunker = None
        self._initialized = False
        self._languages = tesseract_lang or ["eng"]
        self._tessdata_dir = tessdata_path
        self._force_full_page_ocr = force_full_page_ocr
    
    def _create_converter(self, config: ProcessingConfig, formula_enrichment: Optional[bool] = None) -> DocumentConverter:
        """Create a DocumentConverter with Tesseract CLI OCR."""
        # Configure pipeline
        pipeline_options = PdfPipelineOptions()
        pipeline_options.do_ocr = config.do_ocr
        pipeline_options.do_table_structure = config.do_table_structure
        
        # Determine formula enrichment setting (independent of do_ocr)
        if formula_enrichment is not None:
             pipeline_options.do_formula_enrichment = formula_enrichment
        elif not config.formula_ocr:
             # Formula OCR explicitly disabled
             pipeline_options.do_formula_enrichment = False
        elif config.formula_mode == "full":
             pipeline_options.do_formula_enrichment = True
        else:
             # Default to False for "fast" and "smart" (initial pass)
             pipeline_options.do_formula_enrichment = False
        
        # Enable image export
        pipeline_options.generate_page_images = True
        pipeline_options.generate_picture_images = config.export_images
        pipeline_options.images_scale = 2.0
        
        # Use Tesseract CLI for OCR
        ocr_options = TesseractCliOcrOptions(
            lang=self._languages,
            tesseract_cmd="tesseract",
            path=self._tessdata_dir,
            force_full_page_ocr=config.force_full_page_ocr,
        )
        pipeline_options.ocr_options = ocr_options
        
        return DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
                InputFormat.DOCX: WordFormatOption(pipeline_options=pipeline_options),
                InputFormat.PPTX: PowerpointFormatOption(pipeline_options=pipeline_options),
                InputFormat.XLSX: ExcelFormatOption(pipeline_options=pipeline_options),
                InputFormat.CSV: CsvFormatOption(pipeline_options=pipeline_options),
            }
        )


    def _run_docling(self, file_path: Path, config: ProcessingConfig):
        """Run Docling conversion and return the DoclingDocument."""
        # Check if we need to re-initialize converter due to config change or first run
        # For simplicity, we just ensure self._converter exists. 
        # In smart/fast mode, it will be the "fast" converter (no enrichment).
        # In full mode, it will be the "full" converter (enrichment).
        if not self._initialized:
            logger.info("Initializing Docling pipeline...")
            self._converter = self._create_converter(config)
            self._chunker = HierarchicalChunker()
            self._initialized = True
            logger.info(f"Docling pipeline initialized (formula_mode={config.formula_mode})")

        file_path = Path(file_path)
        
        logger.info(f"Extracting with Docling: {file_path.name}")
        
        try:
            # Powerpoint/Excel/Word/CSV use standard conversion
            ext = file_path.suffix.lower()
            if ext not in [".pdf"]:
                result = self._converter.convert(str(file_path))
                
                # DOCX/PPTX: inject OMML equations as LaTeX
                if ext in (".docx", ".pptx") and config.formula_mode != "fast":
                    if ext == ".docx":
                        latex_eqs = self._extract_docx_equations(file_path)
                    else:
                        latex_eqs = self._extract_pptx_equations(file_path)
                    
                    if latex_eqs:
                        # Find formula blocks in Docling output (order-based)
                        formula_blocks = []
                        for item, _ in result.document.iterate_items():
                            label = getattr(item, "label", None)
                            if label and ("formula" in str(label).lower() or "equation" in str(label).lower()):
                                formula_blocks.append(item)
                        
                        # Order-based substitution with alignment gate
                        injected = 0
                        _non_omml = 0
                        for block, latex in zip(formula_blocks, latex_eqs):
                            orig_len = len(block.text.strip()) if block.text else 0
                            latex_len = len(latex.strip())
                            
                            # Asymmetric gate: allow if Docling text is empty/garbled
                            if orig_len < 3 and latex_len > 3:
                                block.text = f"$${latex}$$"
                                injected += 1
                            elif latex_len > 0 and 0.2 <= (orig_len + 5) / (latex_len + 5) <= 5.0:
                                block.text = f"$${latex}$$"
                                injected += 1
                            else:
                                logger.debug(f"Skipping equation inject: ratio out of range")
                        
                        if len(formula_blocks) != len(latex_eqs):
                            logger.warning(
                                f"{ext.upper()} equation count mismatch: "
                                f"extracted={len(latex_eqs)}, docling={len(formula_blocks)}. "
                                f"Injected {injected}."
                            )
                        else:
                            logger.info(f"Injected {injected} LaTeX equations from {ext.upper()}")
                
                return result

            # --- PDF Handling with Smart Formula Mode ---
            
            # Pass 1: Run standard conversion
            # If mode="smart" or "fast", this is the FAST pass (no enrichment).
            # If mode="full", this is the FULL pass (enrichment enabled).
            start_time = time.time()
            result = self._converter.convert(str(file_path))
            _duration = time.time() - start_time
            
            # If not smart mode, we are done (Full or Fast)
            if config.formula_mode != "smart":
                # Apply normalization if Fast mode (to make unicode math nicer)
                if config.formula_mode == "fast":
                    _keys = list(result.document.pages.keys())  # snapshot (unused; iteration below)
                    for _, page in result.document.pages.items():
                         # Iterate all items on page
                         # We can't easily modify text in-place efficiently without iterating items
                         pass 
                    # Actually, normalization is better applied globally to the doc text items
                    for item, _ in result.document.iterate_items():
                        if hasattr(item, "text"):
                            item.text = self._normalize_unicode_math(item.text)
                return result

            # --- Smart Mode: BBox Crop → LaTeX-OCR ---

            num_pages = len(result.document.pages)
            
            # Page cap: if PDF is huge, fallback to fast
            if num_pages > 100:
                logger.info(f"Smart mode disabled: {num_pages} pages exceeds cap (100). "
                            "Falling back to Unicode normalization only.")
                for item, _ in result.document.iterate_items():
                    if hasattr(item, "text"):
                        item.text = self._normalize_unicode_math(item.text)
                return result

            # Find equation items (FORMULA-labeled blocks only)
            equation_items = self._find_equation_items(result.document)
            
            if equation_items:
                # Merge adjacent formula fragments
                merged_items, union_bboxes, blank_ids = self._merge_adjacent_formulas(
                    equation_items, result.document
                )
                
                backend = os.getenv("LONGPARSER_LATEX_OCR_BACKEND", "pix2tex")
                try:
                    from .latex_ocr import LaTeXOCR
                    ocr = LaTeXOCR(backend=backend)
                except ImportError:
                    ocr = None
                    logger.warning("latex_ocr module not available. Skipping formula OCR.")
                
                if ocr and ocr.available:
                    processed, t0 = 0, time.monotonic()
                    # Per-equation timeout: cap each pix2tex call to prevent one slow eq from blocking
                    per_eq_timeout = float(os.getenv("LONGPARSER_FORMULA_PER_EQ_TIMEOUT", "30"))
                    from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeout
                    executor = ThreadPoolExecutor(max_workers=1)
                    
                    for item, page_no in merged_items:
                        # Circuit breaker: equation count
                        if processed >= config.smart_max_equations:
                            logger.info(f"Circuit breaker: {processed} equations reached limit")
                            break
                        # Circuit breaker: time budget
                        if time.monotonic() - t0 > config.smart_max_ocr_seconds:
                            logger.info(f"Circuit breaker: OCR time budget exceeded "
                                        f"({time.monotonic() - t0:.1f}s > {config.smart_max_ocr_seconds}s)")
                            break
                        
                        crop = self._crop_equation_bbox(
                            result.document, item, page_no, union_bboxes
                        )
                        if crop is None:
                            continue
                        
                        # Run pix2tex with per-equation timeout
                        try:
                            future = executor.submit(ocr.recognize, crop)
                            latex = future.result(timeout=per_eq_timeout)
                        except FuturesTimeout:
                            logger.info(f"Equation OCR timed out after {per_eq_timeout}s, skipping")
                            continue
                        except Exception as e:
                            logger.debug(f"Equation OCR error: {e}")
                            continue
                        
                        if latex:
                            item.text = f"$${latex}$$"
                            processed += 1
                    
                    executor.shutdown(wait=False)
                    
                    # Blank leftover fragments (merged items whose text was absorbed)
                    for item, _ in result.document.iterate_items():
                        if id(item) in blank_ids:
                            item.text = ""
                    
                    logger.info(f"Smart mode: OCR'd {processed} equations in "
                                f"{time.monotonic() - t0:.2f}s")
                else:
                    logger.info("LaTeX-OCR not available. Using Unicode normalization only.")
            else:
                logger.info("Smart mode: No FORMULA blocks detected by Docling.")

            # ── MFD fallback: scan candidate pages for missed equations ───────────
            try:
                from .latex_ocr import MFDBackend
                mfd = MFDBackend.get()
            except Exception:
                mfd = None

            if mfd and mfd.available and ocr and ocr.available:
                t0_mfd = time.monotonic()
                max_ocr_secs = config.smart_max_ocr_seconds

                for page_no, page_obj in result.document.pages.items():
                    # Budget gate: skip MFD if <60% of budget remains
                    elapsed = time.monotonic() - t0_mfd
                    if elapsed > max_ocr_secs * 0.4:
                        logger.info("MFD: time budget low, stopping page scan")
                        break

                    # Count Docling formulas on this page for candidate gating
                    docling_formula_count = sum(
                        1 for item, pno in result.document.iterate_items()
                        if pno == page_no and str(getattr(item, "label", "")).lower() in {"formula", "equation"}
                    )

                    # Collect page blocks for gating and replace-first
                    # (These are the Block objects we will be building during extraction;
                    # at this stage we check items directly from Docling result)
                    page_text_items = [
                        item for item, pno in result.document.iterate_items()
                        if pno == page_no and hasattr(item, "text") and item.text
                    ]

                    # Build minimal proxy dicts for _is_mfd_candidate
                    page_proxy = [
                        type("_P", (), {"text": it.text, "type": str(getattr(it, "label", ""))})  # type: ignore
                        for it in page_text_items
                    ]

                    if not _is_mfd_candidate(page_no, page_proxy, docling_formula_count):
                        continue

                    # Get page PIL image (already rendered by Docling in smart mode)
                    page_img = None
                    try:
                        page_img = page_obj.image.pil_image
                    except Exception:
                        continue
                    if page_img is None:
                        continue

                    mfd_boxes = mfd.detect(page_img)
                    if not mfd_boxes:
                        continue

                    # Build pixel-space bboxes for existing Docling FORMULA items on this page
                    existing_formula_px: list[dict] = []
                    img_w, img_h = page_img.size
                    page_w = page_obj.size.width
                    page_h = page_obj.size.height
                    sx = img_w / page_w if page_w else 1.0
                    sy = img_h / page_h if page_h else 1.0

                    for item, pno in result.document.iterate_items():
                        if pno != page_no:
                            continue
                        label = str(getattr(item, "label", "")).lower()
                        if label not in {"formula", "equation"}:
                            continue
                        for prov in getattr(item, "prov", []):
                            if getattr(prov, "page_no", None) != page_no:
                                continue
                            bbox = getattr(prov, "bbox", None)
                            if bbox is None:
                                continue
                            tl = bbox.to_top_left_origin(page_h)
                            existing_formula_px.append({
                                "x0": int(tl.l * sx), "y0": int(tl.t * sy),
                                "x1": int(tl.r * sx), "y1": int(tl.b * sy),
                            })

                    for mbox in mfd_boxes:
                        # Circuit breakers
                        if processed >= config.smart_max_equations:
                            break
                        if time.monotonic() - t0_mfd > max_ocr_secs:
                            break

                        # Skip if already covered by a Docling formula bbox
                        if any(_iou_px(mbox, ex) > 0.5 for ex in existing_formula_px):
                            continue

                        # Crop and OCR
                        from PIL import Image as _PILImage
                        pad_x = (mbox["x1"] - mbox["x0"]) * 0.15
                        pad_y = (mbox["y1"] - mbox["y0"]) * 0.15
                        cx0 = max(0, mbox["x0"] - pad_x)
                        cy0 = max(0, mbox["y0"] - pad_y)
                        cx1 = min(img_w, mbox["x1"] + pad_x)
                        cy1 = min(img_h, mbox["y1"] + pad_y)
                        if (cx1 - cx0) < 64 or (cy1 - cy0) < 64:
                            continue
                        crop = page_img.crop((int(cx0), int(cy0), int(cx1), int(cy1)))

                        latex = ocr.recognize(crop)
                        if not latex:
                            continue

                        processed += 1
                        delim = "$$" if mbox["type"] == "isolated" else "$"
                        latex_text = f"{delim}{latex}{delim}"
                        mbox_dict = mbox  # alias for IoU below

                        # Replace-first: find an overlapping garbled non-formula item
                        replaced = False
                        for item, pno in result.document.iterate_items():
                            if pno != page_no:
                                continue
                            if not hasattr(item, "text") or not item.text:
                                continue
                            label = str(getattr(item, "label", "")).lower()
                            if label in {"formula", "equation"}:
                                continue
                            if not _MATH_RE.search(item.text):
                                continue
                            # Compute pixel bbox for this item
                            for prov in getattr(item, "prov", []):
                                if getattr(prov, "page_no", None) != page_no:
                                    continue
                                bbox = getattr(prov, "bbox", None)
                                if bbox is None:
                                    continue
                                tl = bbox.to_top_left_origin(page_h)
                                item_px = {
                                    "x0": int(tl.l * sx), "y0": int(tl.t * sy),
                                    "x1": int(tl.r * sx), "y1": int(tl.b * sy),
                                }
                                if _iou_px(item_px, mbox_dict) > 0.5:
                                    item.text = latex_text
                                    # Update label to formula so downstream sees it correctly
                                    try:
                                        item.label = type(item.label)("formula")
                                    except Exception:
                                        pass
                                    replaced = True
                                    logger.debug(f"MFD: replaced garbled block on page {page_no}")
                                    break
                            if replaced:
                                break

                        if not replaced:
                            # Append a new synthetic formula item text to the first item on
                            # this page so it flows into the block extraction pass.
                            # Simpler: log and let the extractor create it via block loop.
                            logger.debug(f"MFD: no overlapping garbled block found on page {page_no}; "
                                         f"new equation injected as standalone")
                            # Inject as a minimal TextItem appended to the page's item list
                            try:
                                from docling_core.types.doc import TextItem as _TextItem, DocItemLabel as _DIL
                                new_item = _TextItem(
                                    label=_DIL.FORMULA,
                                    text=latex_text,
                                    prov=[],
                                )
                                result.document.texts.append(new_item)
                            except Exception as e:
                                logger.debug(f"MFD: could not inject new item: {e}")

                logger.info(f"MFD fallback finished. Total OCR'd: {processed}")
            
            # Normalize remaining text (items not replaced with LaTeX)
            for item, _ in result.document.iterate_items():
                if hasattr(item, "text") and not item.text.startswith("$$"):
                    item.text = self._normalize_unicode_math(item.text)

            return result
            
        except Exception as e:
            logger.error(f"Docling extraction failed: {e}")
            raise

    def _cluster_font_sizes(self, heights: List[float], tolerance: float = 0.15) -> List[List[float]]:
        """
        Cluster heading bbox heights into distinct font-size groups.
        
        Uses relative tolerance: two heights belong to the same cluster
        if they are within `tolerance` (15%) of the cluster's mean height.
        
        Returns:
            List of clusters, sorted from largest mean height to smallest.
            Each cluster is a list of heights that belong to it.
        """
        if not heights:
            return []
        
        sorted_heights = sorted(set(heights), reverse=True)
        clusters = []
        
        for h in sorted_heights:
            placed = False
            for cluster in clusters:
                cluster_mean = sum(cluster) / len(cluster)
                # Relative difference check
                if abs(h - cluster_mean) / max(cluster_mean, 0.1) <= tolerance:
                    cluster.append(h)
                    placed = True
                    break
            if not placed:
                clusters.append([h])
        
        # Sort clusters by mean height descending (largest font first)
        clusters.sort(key=lambda c: sum(c) / len(c), reverse=True)
        return clusters

    @staticmethod
    def _extract_marker(text: str) -> Optional[str]:
        """
        Extract the leading marker/prefix from a heading text.
        
        Detects structured prefixes like "I.", "A.", "1.", "IV.",
        "2.3", etc. using a general pattern matcher.
        
        Returns:
            The marker string if found, or None.
        """
        m = _MARKER_RE.match(text.strip())
        return m.group(1) if m else None

    @staticmethod
    def _classify_marker_type(marker: str) -> str:
        """
        Classify a marker using strict numbering patterns.
        
        Returns:
            'numeric' for 1, 1.1, 1.1.1
            'alpha' for A, B, A.1
            'roman' for I, II, IV
            'other' for bullets or non-structural markers
        """
        if not marker:
            return 'other'
        
        marker = marker.strip()
        
        # Strict Numeric: 1. or 1.1 or 1.1.1
        if re.match(r'^\d+(\.\d+)*\.?$', marker):
            return 'numeric'
        
        # Strict Roman: I. or IV. (common uppercase roman)
        if re.match(r'^[IVX]+\.?$', marker):
            return 'roman'
            
        # Strict Alpha: A. or A.1
        if re.match(r'^[A-Z](\.[0-9]+)*\.?$', marker):
            return 'alpha'
            
        return 'other'

    def _sub_cluster_by_markers(
        self,
        texts_in_cluster: List[str],
        base_level: int,
    ) -> Dict[str, int]:
        """
        Sub-differentiate headings within the same font-size cluster
        using autonomous marker-pattern analysis.
        
        Fully data-driven — no hardcoded rankings or character sets.
        
        Algorithm:
        1. Group headings by marker character-class (objective string
           properties: isdigit, isupper, len).
        2. Compute average span (gap) between consecutive markers of
           each group. Parent sections have LARGER spans because child
           sections fill the gaps between them.
        3. Rank groups by span size: largest span = parent level.
        
        Args:
            texts_in_cluster: Heading texts in this font-size cluster.
            base_level: Level assigned by font-size clustering.
            
        Returns:
            Dict mapping heading text -> adjusted heading level.
        """
        if len(texts_in_cluster) <= 1:
            return {t: base_level for t in texts_in_cluster}
        
        # Step 1: Extract and classify markers by character class
        text_info = []  # [(text, mtype)]
        type_counts = {}
        
        for text in texts_in_cluster:
            marker = self._extract_marker(text)
            mtype = self._classify_marker_type(marker) if marker else None
            text_info.append((text, mtype))
            if mtype:
                type_counts[mtype] = type_counts.get(mtype, 0) + 1
        
        # Need at least 2 distinct types with 2+ headings each
        active_types = {t for t, c in type_counts.items() if c >= 2}
        if len(active_types) <= 1:
            return {t: base_level for t in texts_in_cluster}
        
        # Step 2: Compute average span for each marker type
        # Parent groups have LARGER spans (children fill the gaps)
        type_positions = {}
        for idx, (text, mtype) in enumerate(text_info):
            if mtype and mtype in active_types:
                type_positions.setdefault(mtype, []).append(idx)
        
        type_avg_span = {}
        for mtype, positions in type_positions.items():
            if len(positions) < 2:
                # Single instance — treat as broadest span
                type_avg_span[mtype] = len(text_info)
            else:
                spans = [positions[i+1] - positions[i]
                         for i in range(len(positions) - 1)]
                type_avg_span[mtype] = sum(spans) / len(spans)
        
        # Step 3: Sort by average span DESCENDING (largest = parent)
        sorted_types = sorted(
            active_types,
            key=lambda t: type_avg_span[t],
            reverse=True,
        )
        
        # Assign sub-levels
        type_to_sublevel = {}
        for i, mtype in enumerate(sorted_types):
            type_to_sublevel[mtype] = base_level + i
        
        result = {}
        for text, mtype in text_info:
            if mtype and mtype in type_to_sublevel:
                result[text] = type_to_sublevel[mtype]
            else:
                result[text] = base_level
        
        logger.debug(
            f"Sub-clustered {len(texts_in_cluster)} headings at level {base_level}: "
            f"types={dict(type_counts)}, spans={type_avg_span}, "
            f"sub-levels={type_to_sublevel}"
        )
        
        return result

    def _build_hierarchy_map(self, docling_doc) -> Tuple[Dict[str, List[str]], Dict[str, int]]:
        """
        Build two mappings using Docling's native APIs:
        1. item self_ref -> heading path (from HierarchicalChunker)
        2. heading text -> heading level (font-size + marker analysis)
        
        Two-phase heading level inference:
          Phase 1: Font-size clustering — groups by bbox height.
                   Largest font = h1, next = h2, etc.
          Phase 2: Marker-pattern analysis — within each font-size
                   cluster, detect structural prefix patterns to
                   create sub-levels (e.g. "I." parent, "A." child).
        
        Returns:
            Tuple of (ref_to_path, heading_to_level)
        """
        ref_to_path = {}
        heading_to_level = {}
        
        # --- Step 1: Collect heading texts and bbox heights ---
        heading_heights = {}  # text -> height
        heading_order = []    # preserve document order
        
        for item, level in docling_doc.iterate_items():
            if isinstance(item, SectionHeaderItem):
                text = getattr(item, 'text', '')
                if not text:
                    continue
                
                height = 0.0
                prov = getattr(item, 'prov', [])
                if prov and len(prov) > 0:
                    bbox = getattr(prov[0], 'bbox', None)
                    if bbox:
                        height = abs(getattr(bbox, 't', 0) - getattr(bbox, 'b', 0))
                
                if text not in heading_heights:
                    heading_heights[text] = height
                    heading_order.append(text)
        
        if not heading_heights:
            logger.info("No section headers found in document")
            return ref_to_path, heading_to_level
        
        # --- Step 2: Font-size clustering ---
        all_heights = list(heading_heights.values())
        clusters = self._cluster_font_sizes(all_heights)
        
        # Build height -> cluster index
        height_to_cidx = {}
        for idx, cluster in enumerate(clusters):
            for h in cluster:
                height_to_cidx[h] = idx
        
        # Group heading texts by font-size cluster
        cluster_texts = {}
        for text in heading_order:
            cidx = height_to_cidx.get(heading_heights[text], 0)
            if cidx not in cluster_texts:
                cluster_texts[cidx] = []
            cluster_texts[cidx].append(text)
        
        # --- Step 3: Marker-pattern sub-clustering & Late-Arrival Logic ---
        
        # 3a. Find the first "Strong" (Numbered) heading in the entire document
        first_strong_index = float('inf')
        for idx, text in enumerate(heading_order):
            marker = self._extract_marker(text)
            mtype = self._classify_marker_type(marker) if marker else 'other'
            if mtype in ('numeric', 'alpha', 'roman'):
                first_strong_index = idx
                break
        
        logger.info(f"First strong heading index: {first_strong_index if first_strong_index != float('inf') else 'None'}")
        
        # 3b. Assign levels with Late-Arrival check
        current_level = 1
        
        for cidx in sorted(cluster_texts.keys()):
            texts = cluster_texts[cidx]
            
            # Filter matches for this cluster
            valid_texts = []
            demoted_texts = []
            
            for text in texts:
                # Global index in the document
                g_idx = heading_order.index(text)
                
                marker = self._extract_marker(text)
                mtype = self._classify_marker_type(marker) if marker else 'other'
                
                # Late Arrival Rule:
                # If Unnumbered AND appears AFTER the first strong heading -> Demote
                if mtype == 'other':
                    # Allow standard titles even if unnumbered
                    is_standard = text.strip().lower() in {
                        "introduction", "abstract", "background", "objective", 
                        "conclusion", "references", "appendix"
                    }
                    if not is_standard and g_idx > first_strong_index:
                         # Demote to -1 (Paragraph) or a very deep level?
                         # Decision: Demote to -1 to force Paragraph type
                         demoted_texts.append(text)
                         continue
                
                valid_texts.append(text)
            
            # Apply levels to valid texts
            if valid_texts:
                sub_levels = self._sub_cluster_by_markers(valid_texts, base_level=current_level)
                heading_to_level.update(sub_levels)
                max_sub = max(sub_levels.values()) if sub_levels else current_level
                current_level = max_sub + 1
            
            # Apply demotion (-1 -> Paragraph)
            for t in demoted_texts:
                heading_to_level[t] = -1
                
        # Log results
        level_counts = {}
        demoted_count = 0
        for lvl in heading_to_level.values():
            if lvl == -1:
                demoted_count += 1
            else:
                level_counts[lvl] = level_counts.get(lvl, 0) + 1
        
        cluster_info = ", ".join(
            f"h{i+1}={sum(c)/len(c):.1f}px ({len(c)} headings)"
            for i, c in enumerate(clusters)
        )
        logger.info(
            f"Heading levels analyzed: {len(heading_to_level)} total. "
            f"Valid levels={dict(sorted(level_counts.items()))}, "
            f"Demoted (Text)={demoted_count} "
            f"[clusters: {cluster_info}]"
        )
        
        # --- Step 4: Build ref_to_path from HierarchicalChunker ---
        try:
            chunks = list(self._chunker.chunk(docling_doc))
            for chunk in chunks:
                heading_path = []
                if hasattr(chunk, 'meta') and chunk.meta:
                    if hasattr(chunk.meta, 'headings') and chunk.meta.headings:
                        heading_path = list(chunk.meta.headings)
                    
                    if hasattr(chunk.meta, 'doc_items') and chunk.meta.doc_items:
                        for item in chunk.meta.doc_items:
                            ref = getattr(item, 'self_ref', None)
                            if ref:
                                ref_to_path[ref] = heading_path
        except Exception as e:
            logger.warning(f"HierarchicalChunker failed, hierarchy paths will be empty: {e}")
        
        return ref_to_path, heading_to_level

    def _get_page_dimensions(self, docling_doc) -> Dict[int, Tuple[float, float]]:
        """
        Extract actual page dimensions from Docling document.
        
        Returns:
            Dict mapping page_no (0-based) -> (width, height)
        """
        dims = {}
        if hasattr(docling_doc, 'pages') and docling_doc.pages:
            for page_no, page in docling_doc.pages.items():
                width, height = 612.0, 792.0  # Fallback to US Letter
                if hasattr(page, 'size') and page.size:
                    width = float(page.size.width) if hasattr(page.size, 'width') else 612.0
                    height = float(page.size.height) if hasattr(page.size, 'height') else 792.0
                dims[page_no - 1] = (width, height)  # Convert to 0-based
        return dims

    def _extract_bbox(self, prov) -> BoundingBox:
        """Extract BoundingBox from a provenance entry."""
        if not prov or not hasattr(prov, 'bbox') or not prov.bbox:
            return BoundingBox(x0=0, y0=0, x1=0, y1=0)
        
        prov_bbox = prov.bbox
        if hasattr(prov_bbox, 'l'):
            return BoundingBox(
                x0=float(prov_bbox.l),
                y0=float(prov_bbox.t),
                x1=float(prov_bbox.r),
                y1=float(prov_bbox.b),
            )
        elif isinstance(prov_bbox, (list, tuple)) and len(prov_bbox) >= 4:
            return BoundingBox(
                x0=float(prov_bbox[0]),
                y0=float(prov_bbox[1]),
                x1=float(prov_bbox[2]),
                y1=float(prov_bbox[3]),
            )
        return BoundingBox(x0=0, y0=0, x1=0, y1=0)

    def _get_item_provenance(self, item) -> Tuple[int, BoundingBox]:
        """
        Extract page number (0-based) and bbox from a Docling item.
        
        Returns:
            Tuple of (page_number_0based, BoundingBox)
        """
        page_num = 0
        bbox = BoundingBox(x0=0, y0=0, x1=0, y1=0)
        
        if hasattr(item, 'prov') and item.prov:
            for prov in item.prov:
                if hasattr(prov, 'page_no'):
                    page_num = prov.page_no - 1  # Convert to 0-based
                bbox = self._extract_bbox(prov)
                break  # Use first provenance entry
        
        return page_num, bbox

    def _determine_block_type(self, item, level: int, heading_to_level: Dict[str, int] = None) -> Tuple[BlockType, Optional[int]]:
        """
        Determine block type and heading level from a Docling item
        using isinstance checks and item.label.
        
        For headings, uses the heading_to_level map (built from
        HierarchicalChunker) for proper heading depth.
        
        Returns:
            Tuple of (BlockType, heading_level_or_None)
        """
        heading_level = None
        
        # Primary: isinstance checks (most reliable)
        if isinstance(item, SectionHeaderItem):
            # Use HierarchicalChunker-derived level if available
            text = getattr(item, 'text', '')
            if heading_to_level and text in heading_to_level:
                lvl = heading_to_level[text]
                # If level is -1, it was demoted to Paragraph
                if lvl == -1:
                    return BlockType.PARAGRAPH, None
                heading_level = lvl
            else:
                heading_level = max(1, level)
            return BlockType.HEADING, heading_level
        
        # PPTX slide titles come as TitleItem (extends TextItem)
        if TitleItem is not None and isinstance(item, TitleItem):
            text = getattr(item, 'text', '')
            if heading_to_level and text in heading_to_level:
                lvl = heading_to_level[text]
                if lvl == -1:
                    return BlockType.PARAGRAPH, None
                heading_level = lvl
            else:
                heading_level = max(1, level)
            return BlockType.HEADING, heading_level
        
        if isinstance(item, TableItem):
            return BlockType.TABLE, None
        
        if isinstance(item, ListItem):
            return BlockType.LIST_ITEM, None
        
        if isinstance(item, PictureItem):
            return BlockType.FIGURE, None
        
        # Secondary: check item.label for fine-grained classification
        label = getattr(item, 'label', None)
        if label:
            label_str = str(label).lower() if not isinstance(label, str) else label.lower()
            
            if 'caption' in label_str:
                return BlockType.CAPTION, None
            if 'footer' in label_str or 'footnote' in label_str:
                return BlockType.FOOTER, None
            if 'header' in label_str and 'section' not in label_str:
                return BlockType.HEADER, None
            if 'equation' in label_str or 'formula' in label_str:
                return BlockType.EQUATION, None
            if 'code' in label_str:
                return BlockType.CODE, None
            if 'title' in label_str:
                return BlockType.HEADING, max(1, level)
        
        # Default: paragraph
        return BlockType.PARAGRAPH, None

    def _get_item_text(self, item, docling_doc=None) -> str:
        """Extract text from a Docling item."""
        # For tables, prefer markdown with doc context for proper rendering
        if isinstance(item, TableItem) and hasattr(item, 'export_to_markdown'):
            try:
                return item.export_to_markdown(doc=docling_doc)
            except Exception:
                pass
        if hasattr(item, 'text') and item.text:
            return item.text
        if hasattr(item, 'export_to_markdown'):
            try:
                return item.export_to_markdown()
            except Exception:
                pass
        return ""

    def _get_item_confidence(self, item) -> float:
        """Extract confidence from a Docling item, defaulting to 1.0."""
        if hasattr(item, 'confidence') and item.confidence is not None:
            return float(item.confidence)
        return 1.0

    def _build_pptx_text_map(self, file_path: Path) -> Dict[int, Dict[str, PptxParaInfo]]:
        """
        Use python-pptx to build a per-slide map of text -> paragraph info.
        
        Returns:
            Dict[slide_idx (0-based), Dict[normalized_text, PptxParaInfo]]
        """
        try:
            from pptx import Presentation
            from pptx.util import Emu
            from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PP_PLACEHOLDER
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            logger.warning("python-pptx not installed, cannot build PPTX indent map")
            return {}
        
        pptx_map: Dict[int, Dict[str, PptxParaInfo]] = {}
        
        try:
            prs = Presentation(str(file_path))
        except Exception as e:
            logger.warning(f"Failed to open PPTX with python-pptx: {e}")
            return {}
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_map: Dict[str, PptxParaInfo] = {}
            found_title = False
            
            # Check if slide 0 has an actual SUBTITLE placeholder
            # If it does, we don't need the positional heuristic
            has_subtitle_placeholder = False
            if slide_idx == 0:
                try:
                    from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PP_PH
                    for s in slide.shapes:
                        if s.is_placeholder:
                            try:
                                if s.placeholder_format.type == PP_PH.SUBTITLE:
                                    has_subtitle_placeholder = True
                                    break
                            except Exception:
                                pass
                except ImportError:
                    pass
            
            for shape in slide.shapes:
                found_title = self._extract_pptx_shape_info(
                    shape, slide_map, slide_idx=slide_idx, found_title=found_title,
                    has_subtitle_placeholder=has_subtitle_placeholder,
                )
            
            pptx_map[slide_idx] = slide_map
        
        # Post-processing: detect repeated text across slides (footer/header noise)
        # Text appearing on >50% of slides is likely a repeated footer/header element
        num_slides = len(pptx_map)
        if num_slides >= 3:  # Only apply for presentations with enough slides
            text_slide_count: Dict[str, int] = {}
            for slide_map in pptx_map.values():
                for text, info in slide_map.items():
                    if not info.is_title and not info.is_subtitle and not info.is_footer:
                        text_slide_count[text] = text_slide_count.get(text, 0) + 1
            
            threshold = num_slides * 0.5
            repeated_texts = {t for t, count in text_slide_count.items() if count > threshold}
            
            if repeated_texts:
                logger.info(f"Detected {len(repeated_texts)} repeated footer/header texts across slides")
                for slide_map in pptx_map.values():
                    for text in repeated_texts:
                        if text in slide_map:
                            slide_map[text] = PptxParaInfo(
                                indent_level=0, is_title=False, is_subtitle=False,
                                is_list=False, bullet_type='None', is_footer=True,
                            )
        
        total_entries = sum(len(m) for m in pptx_map.values())
        logger.info(f"Built PPTX text map: {len(pptx_map)} slides, {total_entries} text entries")
        return pptx_map
    
    def _extract_pptx_shape_info(self, shape, slide_map: Dict[str, PptxParaInfo], 
                                slide_idx: int = 0, found_title: bool = False,
                                has_subtitle_placeholder: bool = False) -> bool:
        """Extract paragraph info from a shape, handling groups recursively.
        
        Returns whether a title shape has been found (for subtitle detection).
        """
        try:
            from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PP_PLACEHOLDER
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return found_title
        
        # Handle group shapes recursively
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child_shape in shape.shapes:
                found_title = self._extract_pptx_shape_info(
                    child_shape, slide_map, slide_idx=slide_idx, found_title=found_title,
                    has_subtitle_placeholder=has_subtitle_placeholder,
                )
            return found_title
        
        if not hasattr(shape, 'text_frame'):
            return found_title
        if not shape.has_text_frame:
            return found_title
        
        # Determine if this shape is a title/subtitle/footer placeholder
        is_title_shape = False
        is_subtitle_shape = False
        is_footer_shape = False
        if shape.is_placeholder:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    is_title_shape = True
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    is_subtitle_shape = True
                elif ph_type in (PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER):
                    is_footer_shape = True
            except Exception:
                pass
        
        # Skip footer/date/slide-number shapes entirely
        if is_footer_shape:
            # Still record them in the map so we can filter during block conversion
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    norm_text = ' '.join(text.split())
                    if norm_text not in slide_map:
                        slide_map[norm_text] = PptxParaInfo(
                            indent_level=0, is_title=False, is_subtitle=False,
                            is_list=False, bullet_type='None', is_footer=True,
                        )
            return found_title
        
        # Detect subtitle: on the title slide (slide 0), the first non-placeholder
        # text shape after the TITLE is a subtitle — but ONLY if there's no actual
        # SUBTITLE placeholder on the slide (to avoid false positives)
        is_subtitle_by_position = False
        if (slide_idx == 0 and found_title and not is_title_shape 
                and not shape.is_placeholder and not has_subtitle_placeholder):
            is_subtitle_by_position = True
        
        first_para_in_shape = True
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            
            indent_level = paragraph.level if paragraph.level is not None else 0
            
            # Detect bullet/numbered list
            is_list = False
            bullet_type = 'None'
            p_elem = paragraph._element
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            if p_elem.find('.//a:buChar', namespaces=ns) is not None:
                is_list = True
                bullet_type = 'Bullet'
            elif p_elem.find('.//a:buAutoNum', namespaces=ns) is not None:
                is_list = True
                bullet_type = 'Numbered'
            elif indent_level > 0 and not is_title_shape:
                is_list = True
            
            # For title shapes, override indent to 0
            if is_title_shape or is_subtitle_shape:
                indent_level = 0
                is_list = False
                bullet_type = 'None'
            
            # Mark as subtitle only for the FIRST paragraph of the subtitle shape
            mark_subtitle = is_subtitle_shape or (is_subtitle_by_position and first_para_in_shape)
            if mark_subtitle:
                indent_level = 0
                is_list = False
                bullet_type = 'None'
            
            # Normalize text for matching (Docling may strip/normalize differently)
            norm_text = ' '.join(text.split())
            
            # Only store first occurrence per slide (duplicate text on same slide is rare)
            if norm_text not in slide_map:
                slide_map[norm_text] = PptxParaInfo(
                    indent_level=indent_level,
                    is_title=is_title_shape,
                    is_subtitle=mark_subtitle,
                    is_list=is_list,
                    bullet_type=bullet_type,
                )
            
            first_para_in_shape = False
        
        # Track that we've seen a title shape
        if is_title_shape:
            found_title = True
        
        return found_title

    def extract(
        self,
        file_path: Path,
        config: ProcessingConfig,
        page_numbers: Optional[List[int]] = None,
    ) -> Tuple[Document, ExtractionMetadata]:
        """
        Extract document using Docling.
        
        Uses iterate_items() for reading-order block extraction
        and HierarchicalChunker for heading hierarchy paths.
        For PPTX files, uses python-pptx directly for indent levels.
        
        Args:
            file_path: Path to document file
            config: Processing configuration
            page_numbers: Optional specific pages to extract
            
        Returns:
            Tuple of (Document, ExtractionMetadata)
        """
        file_path = Path(file_path)
        is_pptx = file_path.suffix.lower() in ('.pptx', '.ppt')
        
        # Calculate file hash
        with open(file_path, "rb") as f:
            file_hash = hashlib.md5(f.read()).hexdigest()
        
        # Get conversion result (cached or new)
        result = self._run_docling(file_path, config)
        docling_doc = result.document
        
        # Build PPTX-specific indent map if applicable
        pptx_text_map = None
        if is_pptx:
            pptx_text_map = self._build_pptx_text_map(file_path)
            # For PPTX: skip font-size clustering, use simple heading levels
            # All slide titles become h2 (since they're all peer-level slides)
            heading_to_level = {}
            for item, level in docling_doc.iterate_items():
                if isinstance(item, SectionHeaderItem):
                    text = getattr(item, 'text', '')
                    if text:
                        heading_to_level[text] = 2  # All PPTX titles = h2
                elif TitleItem is not None and isinstance(item, TitleItem):
                    text = getattr(item, 'text', '')
                    if text:
                        heading_to_level[text] = 2  # All PPTX titles = h2
            hierarchy_map = {}
            logger.info(f"PPTX mode: assigned {len(heading_to_level)} headings to level 2")
        else:
            # Standard PDF/DOCX path: use font-size clustering
            hierarchy_map, heading_to_level = self._build_hierarchy_map(docling_doc)
        
        logger.info(f"Built hierarchy map with {len(hierarchy_map)} item mappings")
        
        # Get actual page dimensions
        page_dims = self._get_page_dimensions(docling_doc)
        
        # Convert to our Document format using iterate_items()
        pages = self._convert_to_pages(
            docling_doc, 
            hierarchy_map, 
            heading_to_level, 
            page_dims, 
            file_path, 
            file_hash,
            exclude_headers_footers=config.exclude_page_headers_footers,
            pptx_text_map=pptx_text_map,
        )
        
        # Filter pages if specific ones requested
        if page_numbers is not None:
            pages = [p for p in pages if p.page_number in page_numbers]
        
        # Build document
        doc = Document(
            metadata=DocumentMetadata(
                source_file=str(file_path),
                file_hash=file_hash,
                total_pages=len(pages),
            ),
            pages=pages,
        )
        
        # Extraction metadata
        strategy_desc = "PPTX mode (python-pptx indent map)" if is_pptx else "PDF/DOCX mode (font-size clustering)"
        meta = ExtractionMetadata(
            strategy_used="docling",
            ocr_backend_used="tesseract_cli",
            reasons=[f"Used Docling with Tesseract CLI OCR, iterate_items() for block extraction. {strategy_desc}"],
        )
        
        total_blocks = sum(len(p.blocks) for p in pages)
        logger.info(f"Extracted {len(pages)} pages, {total_blocks} blocks")
        
        return doc, meta
    
    def _build_table_from_item(self, item, docling_doc=None) -> Optional[Table]:
        """
        Convert Docling TableItem.data into our Table schema.
        
        Uses table_cells with proper row/col indices.
        Falls back to export_to_dataframe() if direct conversion fails.
        """
        if not isinstance(item, TableItem) or not hasattr(item, 'data'):
            return None
        
        table_data = item.data
        if table_data.num_rows == 0 or table_data.num_cols == 0:
            return None
        
        try:
            cells = []
            for dcell in table_data.table_cells:
                cells.append(TableCell(
                    r0=dcell.start_row_offset_idx,
                    c0=dcell.start_col_offset_idx,
                    rspan=dcell.end_row_offset_idx - dcell.start_row_offset_idx,
                    cspan=dcell.end_col_offset_idx - dcell.start_col_offset_idx,
                    text=dcell.text,
                ))
            
            if cells:
                return Table(
                    n_rows=table_data.num_rows,
                    n_cols=table_data.num_cols,
                    cells=cells,
                )
        except Exception as e:
            logger.warning(f"Direct table cell conversion failed: {e}")
        
        # Fallback: use export_to_dataframe()
        try:
            import pandas as pd
            df = item.export_to_dataframe(doc=docling_doc)
            if df is not None and not df.empty:
                n_rows = len(df) + 1  # +1 for header row
                n_cols = len(df.columns)
                cells = []
                # Header row
                for c_idx, col_name in enumerate(df.columns):
                    cells.append(TableCell(
                        r0=0, c0=c_idx, rspan=1, cspan=1,
                        text=str(col_name),
                    ))
                # Data rows
                for r_idx, (_, row) in enumerate(df.iterrows(), start=1):
                    for c_idx, val in enumerate(row):
                        cells.append(TableCell(
                            r0=r_idx, c0=c_idx, rspan=1, cspan=1,
                            text=str(val) if pd.notna(val) else "",
                        ))
                return Table(
                    n_rows=n_rows,
                    n_cols=n_cols,
                    cells=cells,
                )
        except Exception as e:
            logger.warning(f"DataFrame fallback also failed: {e}")
        
        return None

    def _convert_to_pages(
        self,
        docling_doc,
        hierarchy_map: Dict[str, List[str]],
        heading_to_level: Dict[str, int],
        page_dims: Dict[int, Tuple[float, float]],
        file_path: Path,
        file_hash: str,
        exclude_headers_footers: bool = True,
        pptx_text_map: Optional[Dict[int, Dict[str, 'PptxParaInfo']]] = None,
    ) -> List[Page]:
        """
        Convert Docling document to our Page format using iterate_items().
        
        No synthetic heading injection, no inline heading regex,
        no hardcoded dimensions — purely Docling-native.
        
        Tracks TableItem children to prevent duplicate blocks.
        When pptx_text_map is provided, uses it to set indent_level on blocks.
        """
        pages_dict: Dict[int, Page] = {}
        block_idx = 0
        
        # Gap #1: Collect all self_refs that belong to table children
        # so we can skip them when they appear as standalone items.
        table_child_refs: set = set()
        for item, _level in docling_doc.iterate_items():
            if isinstance(item, TableItem):
                # Mark all refs inside this table's cells as children
                if hasattr(item, 'data') and item.data:
                    for dcell in item.data.table_cells:
                        if hasattr(dcell, 'ref') and dcell.ref:
                            ref = getattr(dcell.ref, 'cref', getattr(dcell.ref, 'self_ref', None))
                            if ref:
                                table_child_refs.add(ref)
        
        # iterate_items() provides (item, level) in reading order
        for item, level in docling_doc.iterate_items():
            # Gap #1: Skip items that are children of a table
            item_ref = getattr(item, 'self_ref', None)
            if item_ref and item_ref in table_child_refs:
                continue
            
            # Get page and bbox from provenance
            page_num, bbox = self._get_item_provenance(item)
            
            # Determine block type using Docling's native types + chunker heading levels
            block_type, heading_level = self._determine_block_type(item, level, heading_to_level)
            
            # Filter headers and footers if requested
            if exclude_headers_footers and block_type in (BlockType.HEADER, BlockType.FOOTER):
                continue
            
            # Get text (prefer markdown for equations to get LaTeX)
            if block_type == BlockType.EQUATION and hasattr(item, 'export_to_markdown'):
                try:
                    text = item.export_to_markdown()
                except Exception:
                    text = self._get_item_text(item, docling_doc)
            else:
                text = self._get_item_text(item, docling_doc)

            if not text:
                continue

            # Wrap equations with markers
            if block_type == BlockType.EQUATION:
                text = f"⟦EQUATION⟧\n{text.strip()}\n⟦/EQUATION⟧"
            
            # Get hierarchy path from chunker map
            item_ref = getattr(item, 'self_ref', None)
            hierarchy_path = hierarchy_map.get(item_ref, [])
            
            # Get native confidence
            item_confidence = self._get_item_confidence(item)
            
            # Create page if needed, with actual dimensions
            if page_num not in pages_dict:
                width, height = page_dims.get(page_num, (612.0, 792.0))
                pages_dict[page_num] = Page(
                    page_number=page_num,
                    width=width,
                    height=height,
                    blocks=[],
                    profile=PageProfile(page_number=page_num),
                )
            
            # Build block
            table_obj = None
            if block_type == BlockType.TABLE:
                table_obj = self._build_table_from_item(item, docling_doc)
                if table_obj:
                    logger.info(
                        f"  Populated Block.table: {table_obj.n_rows} rows × "
                        f"{table_obj.n_cols} cols, {len(table_obj.cells)} cells"
                    )
            
            # Determine indent_level and filter footers from PPTX text map
            indent_level = 0
            if pptx_text_map is not None:
                norm_block_text = ' '.join(text.strip().split())
                # page_num is 1-based, pptx_text_map is 0-based
                slide_map = pptx_text_map.get(page_num - 1, {})
                # Also try page_num as-is (in case of off-by-one)
                if not slide_map:
                    slide_map = pptx_text_map.get(page_num, {})
                pptx_info = slide_map.get(norm_block_text)
                if pptx_info:
                    if pptx_info.is_footer:
                        # Skip footer/date/slide-number content
                        continue
                    indent_level = pptx_info.indent_level
                    # Promote subtitle to heading level 3
                    if pptx_info.is_subtitle and block_type != BlockType.HEADING:
                        block_type = BlockType.HEADING
                        heading_level = 3
                
                # Filter slide number patterns (e.g., "1 / 22", "12/22")
                if re.match(r'^\d+\s*/\s*\d+$', norm_block_text):
                    continue
                # Filter single-character noise (common Beamer artifact)
                if len(norm_block_text) <= 1 and block_type not in (BlockType.HEADING,):
                    continue
            
            block = Block(
                type=block_type,
                text=text,
                order_index=block_idx,
                heading_level=heading_level,
                indent_level=indent_level,
                hierarchy_path=hierarchy_path,
                provenance=Provenance(
                    source_file=str(file_path),
                    page_number=page_num,
                    bbox=bbox,
                    extractor=ExtractorType.DOCLING,
                    extractor_version=self.version,
                ),
                confidence=Confidence(
                    overall=item_confidence,
                    text_confidence=item_confidence,
                    layout_confidence=item_confidence,
                ),
                table=table_obj,
            )
            
            pages_dict[page_num].blocks.append(block)
            block_idx += 1
        
        # Sort pages by page number and reindex blocks
        pages = sorted(pages_dict.values(), key=lambda p: p.page_number)
        for page in pages:
            for i, block in enumerate(page.blocks):
                block.order_index = i
        
        return pages
    
    def extract_page(
        self,
        file_path: Path,
        page_number: int,
        config: ProcessingConfig,
    ) -> Page:
        """Extract a single page."""
        doc, _ = self.extract(file_path, config, page_numbers=[page_number])
        if doc.pages:
            return doc.pages[0]
        raise ValueError(f"Page {page_number} not found in {file_path}")
    
    def get_hierarchy(
        self,
        file_path: Path,
        config: ProcessingConfig,
    ) -> List[HierarchyChunk]:
        """
        Get document hierarchy using HierarchicalChunker.
        
        Returns list of chunks with hierarchy information.
        """
        file_path = Path(file_path)
        
        # Get conversion result (cached or new)
        result = self._run_docling(file_path, config)
        chunks = list(self._chunker.chunk(result.document))
        
        hierarchy_chunks = []
        for idx, chunk in enumerate(chunks):
            heading_path = []
            page_num = 0
            
            if hasattr(chunk, 'meta') and chunk.meta:
                if hasattr(chunk.meta, 'headings'):
                    heading_path = list(chunk.meta.headings or [])
                if hasattr(chunk.meta, 'doc_items') and chunk.meta.doc_items:
                    for item in chunk.meta.doc_items:
                        if hasattr(item, 'prov') and item.prov:
                            for prov in item.prov:
                                if hasattr(prov, 'page_no'):
                                    page_num = prov.page_no - 1
                                    break
            
            hierarchy_chunks.append(HierarchyChunk(
                text=chunk.text,
                heading_path=heading_path,
                level=len(heading_path),
                page_number=page_num,
                order_index=idx,
            ))
        
        return hierarchy_chunks
    
    def to_markdown(self, doc: Document) -> str:
        """Convert document to Markdown."""
        lines = []
        
        for page in doc.pages:
            for block in page.blocks:
                if block.type == BlockType.HEADING and block.heading_level:
                    prefix = "#" * min(block.heading_level, 6)
                    lines.append(f"{prefix} {block.text}")
                    lines.append("")
                elif block.type == BlockType.LIST_ITEM:
                    indent = "  " * block.indent_level
                    lines.append(f"{indent}- {block.text}")
                    lines.append("")
                elif block.type == BlockType.HEADER:
                    # Page headers (e.g. running headers)
                    lines.append(block.text)
                    lines.append("")
                else:
                    text = block.text
                    # Escape leading # in non-heading text to prevent
                    # markdown interpreting code comments as headings
                    if text.lstrip().startswith('#') and block.type != BlockType.HEADING:
                        text = text.replace('#', '\\#', 1)
                    lines.append(text)
                    lines.append("")
        
        return "\n".join(lines)

    def _sanitize_filename(self, name: str) -> str:
        """Sanitize string for filename."""
        return "".join(c for c in name if c.isalnum() or c in ('-', '_')).strip()

    def save_images(self, output_dir: Path) -> List[Path]:
        """
        Save extracted images (pages, figures, tables).
        
        Args:
            output_dir: Directory to save images
            
        Returns:
            List of saved image paths
        """
        if self._last_result:
            result = self._last_result
        else:
            logger.warning("No conversion result available to save images from")
            return []
            
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        saved_paths = []
        
        try:
            # Save page images
            if hasattr(result.document, 'pages'):
                for page_no, page in result.document.pages.items():
                    if hasattr(page, 'image') and page.image and hasattr(page.image, 'pil_image'):
                        image_path = output_dir / f"page_{page_no}.png"
                        try:
                            page.image.pil_image.save(image_path, format="PNG")
                            saved_paths.append(image_path)
                        except Exception as e:
                            logger.warning(f"Failed to save page image {page_no}: {e}")
            
            # Save figures and tables
            for element, _level in result.document.iterate_items():
                if isinstance(element, PictureItem):
                    try:
                        img = element.get_image(result.document)
                        if img:
                            safe_ref = self._sanitize_filename(element.self_ref)
                            if not safe_ref:
                                safe_ref = f"picture_{uuid.uuid4().hex[:8]}"
                            image_path = output_dir / f"figure_{safe_ref}.png"
                            img.save(image_path, format="PNG")
                            saved_paths.append(image_path)
                    except Exception as e:
                        logger.warning(f"Failed to save figure image: {e}")
                
                if isinstance(element, TableItem):
                    try:
                        img = element.get_image(result.document)
                        if img:
                            safe_ref = self._sanitize_filename(element.self_ref)
                            if not safe_ref:
                                safe_ref = f"table_{uuid.uuid4().hex[:8]}"
                            image_path = output_dir / f"table_{safe_ref}.png"
                            img.save(image_path, format="PNG")
                            saved_paths.append(image_path)
                    except Exception as e:
                        logger.warning(f"Failed to save table image: {e}")
                        
        except Exception as e:
            logger.error(f"Failed to save images: {e}")
            
        return saved_paths

    # ------------------------------------------------------------------
    # LaTeX-OCR helpers (PDF smart mode)
    # ------------------------------------------------------------------

    def _find_equation_items(self, doc) -> List[tuple]:
        """Find FORMULA-labeled items. Returns [(item, page_no), ...]."""
        equation_items = []
        for item, _ in doc.iterate_items():
            label = getattr(item, "label", None)
            if label is None:
                continue
            label_str = str(label).lower()
            if "formula" in label_str or "equation" in label_str:
                # Get page number from provenance
                page_no = 1
                if hasattr(item, "prov") and item.prov:
                    page_no = item.prov[0].page_no
                equation_items.append((item, page_no))
        return equation_items

    def _merge_adjacent_formulas(self, items: List[tuple], doc) -> tuple:
        """Merge vertically adjacent FORMULA bboxes in pixel space.

        Returns:
            merged_items: list of (item, page_no)
            union_bboxes: dict of id(item) -> (x0, y0, x1, y1) in pixels
            blank_ids: set of id(item) for leftover fragments to blank
        """
        union_bboxes: Dict[int, tuple] = {}
        blank_ids: set = set()

        if len(items) < 2:
            return items, union_bboxes, blank_ids

        # Convert to pixel-space, matching prov by page_no
        pixel_items = []
        for item, page_no in items:
            # Find provenance matching this page
            prov = None
            for p in getattr(item, 'prov', []):
                if getattr(p, 'page_no', None) == page_no:
                    prov = p
                    break
            if not prov or not prov.bbox:
                continue

            page = doc.pages.get(page_no)
            if page is None or not hasattr(page, 'image') or page.image is None:
                continue

            try:
                pil_img = page.image.pil_image
                img_w, img_h = pil_img.size
                page_w = page.size.width
                page_h = page.size.height

                tl = prov.bbox.to_top_left_origin(page_h)
                sx, sy = img_w / page_w, img_h / page_h
                px_bbox = (tl.l * sx, tl.t * sy, tl.r * sx, tl.b * sy)
                pixel_items.append((item, page_no, px_bbox))
            except Exception as e:
                logger.debug(f"Skipping formula merge for item: {e}")
                continue

        if len(pixel_items) < 2:
            return items, union_bboxes, blank_ids

        GAP_PX = 20
        H_OVERLAP_MIN = 0.3
        OVERLAP_Y_ALLOW = 5

        groups = [[pixel_items[0]]]
        for entry in pixel_items[1:]:
            _, pg, (x0, y0, x1, y1) = entry
            prev = groups[-1][-1]
            _, prev_pg, (px0, py0, px1, py1) = prev

            # Directed gap (y increases downward in pixel space)
            gap = y0 - py1
            if pg != prev_pg or gap < -OVERLAP_Y_ALLOW or gap > GAP_PX:
                groups.append([entry])
                continue

            h_overlap = max(0, min(x1, px1) - max(x0, px0))
            h_extent = max(x1, px1) - min(x0, px0)
            if h_extent > 0 and h_overlap / h_extent >= H_OVERLAP_MIN:
                groups[-1].append(entry)
            else:
                groups.append([entry])

        merged = []
        for group in groups:
            anchor_item, anchor_pg, _ = group[0]
            if len(group) == 1:
                merged.append((anchor_item, anchor_pg))
            else:
                # Compute union bbox in pixel space
                ux0 = min(e[2][0] for e in group)
                uy0 = min(e[2][1] for e in group)
                ux1 = max(e[2][2] for e in group)
                uy1 = max(e[2][3] for e in group)
                union_bboxes[id(anchor_item)] = (ux0, uy0, ux1, uy1)
                merged.append((anchor_item, anchor_pg))
                # Mark non-anchor items for blanking
                for e in group[1:]:
                    blank_ids.add(id(e[0]))

        return merged, union_bboxes, blank_ids

    def _crop_equation_bbox(self, doc, item, page_no: int,
                            union_bboxes: Dict[int, tuple] = None):
        """Crop equation image from page. Returns PIL Image or None."""
        page = doc.pages.get(page_no)
        if page is None or not hasattr(page, 'image') or page.image is None:
            return None

        try:
            pil_img = page.image.pil_image
            img_w, img_h = pil_img.size
        except Exception:
            return None

        # Check for merged union bbox first
        if union_bboxes and id(item) in union_bboxes:
            x0, y0, x1, y1 = union_bboxes[id(item)]
        else:
            # Standard provenance → pixel transform
            prov = None
            for p in getattr(item, 'prov', []):
                if getattr(p, 'page_no', None) == page_no:
                    prov = p
                    break
            if not prov or not prov.bbox:
                return None

            page_w = page.size.width
            page_h = page.size.height

            tl = prov.bbox.to_top_left_origin(page_h)
            sx, sy = img_w / page_w, img_h / page_h
            x0, y0 = tl.l * sx, tl.t * sy
            x1, y1 = tl.r * sx, tl.b * sy

        # Rotation/sanity: coords must be within image
        if x0 < 0 or y0 < 0 or x1 > img_w or y1 > img_h:
            logger.debug(f"BBox outside image bounds on page {page_no}, skipping")
            return None

        # 15% padding + clamp
        pad_x = (x1 - x0) * 0.15
        pad_y = (y1 - y0) * 0.15
        x0, y0 = max(0, x0 - pad_x), max(0, y0 - pad_y)
        x1, y1 = min(img_w, x1 + pad_x), min(img_h, y1 + pad_y)

        # Minimum 64px crop
        if (x1 - x0) < 64 or (y1 - y0) < 64:
            logger.debug(f"Crop too small ({x1-x0:.0f}×{y1-y0:.0f}px) on page {page_no}")
            return None

        return pil_img.crop((int(x0), int(y0), int(x1), int(y1)))

    # ------------------------------------------------------------------
    # DOCX/PPTX equation extraction
    # ------------------------------------------------------------------

    def _extract_docx_equations(self, file_path: Path) -> List[str]:
        """Extract OMML equations from DOCX as LaTeX strings."""
        try:
            from docxlatex import Document as DocxLatexDoc
            doc = DocxLatexDoc(str(file_path))
            equations = doc.get_equations()
            return [self._normalize_latex(eq) for eq in equations if eq.strip()]
        except ImportError:
            logger.warning("docxlatex not installed. Skipping DOCX equation extraction.")
            return []
        except Exception as e:
            logger.warning(f"DOCX equation extraction failed: {e}")
            return []

    def _extract_pptx_equations(self, file_path: Path) -> List[str]:
        """Scan PPTX slide XML for <m:oMath> nodes."""
        import zipfile
        try:
            import defusedxml.ElementTree as ET
        except ImportError:
            logger.warning("defusedxml not installed. Skipping PPTX equation extraction.")
            return []

        equations = []
        MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
        MAX_SLIDES = 50
        MAX_BYTES_PER_SLIDE = 10 * 1024 * 1024
        MAX_TOTAL_BYTES = 100 * 1024 * 1024
        MAX_COMPRESSION_RATIO = 100
        MAX_ENTRIES = 500

        try:
            with zipfile.ZipFile(str(file_path)) as z:
                # Zip entry count cap
                if len(z.infolist()) > MAX_ENTRIES:
                    logger.warning(f"PPTX has {len(z.infolist())} entries (>{MAX_ENTRIES}). Skipping.")
                    return []

                slide_files = sorted([
                    n for n in z.namelist()
                    if n.startswith("ppt/slides/slide") and n.endswith(".xml")
                ])[:MAX_SLIDES]

                total_bytes = 0
                for name in slide_files:
                    info = z.getinfo(name)
                    if info.file_size > MAX_BYTES_PER_SLIDE:
                        logger.debug(f"Skipping {name}: too large")
                        continue
                    if info.compress_size > 0 and info.file_size / info.compress_size > MAX_COMPRESSION_RATIO:
                        logger.debug(f"Skipping {name}: suspicious compression ratio")
                        continue
                    total_bytes += info.file_size
                    if total_bytes > MAX_TOTAL_BYTES:
                        logger.warning("PPTX total uncompressed bytes exceeded cap")
                        break

                    try:
                        tree = ET.parse(z.open(name))
                        for omath in tree.iter(f"{{{MATH_NS}}}oMath"):
                            # Extract text content from m:t elements
                            texts = []
                            for t_elem in omath.iter(f"{{{MATH_NS}}}t"):
                                if t_elem.text:
                                    texts.append(t_elem.text)
                            if texts:
                                raw = " ".join(texts)
                                equations.append(self._normalize_latex(raw))
                    except Exception as e:
                        logger.debug(f"Failed to parse {name}: {e}")
        except Exception as e:
            logger.warning(f"PPTX equation extraction failed: {e}")

        return equations

    def _normalize_latex(self, latex: str) -> str:
        """Fix whitespace artifacts in converted LaTeX."""
        if not latex:
            return latex
        # Collapse broken control sequences: \f r a c → \frac
        prev = None
        while prev != latex:
            prev = latex
            latex = re.sub(r'\\([a-zA-Z]+)\s+([a-zA-Z])', r'\\\1\2', latex)
        # Trim repeated spaces
        latex = re.sub(r'\s+', ' ', latex).strip()
        return latex

    def _normalize_unicode_math(self, text: str) -> str:
        """
        Convert Unicode math symbols to LaTeX-lite notation.
        Only applies if text does NOT look like it already has LaTeX formatting.
        """
        if not text:
            return text
            
        # Scope check: Don't touch if it looks like LaTeX already
        if "$" in text or "\\" in text:
            return text
            
        # Common math symbols
        replacements = [
            (r"²", "^2"), (r"³", "^3"),
            (r"₁", "_1"), (r"₂", "_2"), (r"₃", "_3"), (r"ᵢ", "_i"), (r"ⱼ", "_j"), (r"ₙ", "_n"),
            (r"∑", r"\\sum"), (r"∫", r"\\int"), (r"∞", r"\\infty"),
            (r"√", r"\\sqrt"), (r"∂", r"\\partial"), (r"∇", r"\\nabla"),
            (r"≈", r"\\approx"), (r"≠", r"\\neq"), (r"≤", r"\\leq"), (r"≥", r"\\geq"),
            (r"α", r"\\alpha"), (r"β", r"\\beta"), (r"γ", r"\\gamma"), (r"θ", r"\\theta"),
            (r"π", r"\\pi"), (r"µ", r"\\mu"), (r"σ", r"\\sigma"), (r"Ω", r"\\Omega"),
            (r"∈", r"\\in"), (r"∀", r"\\forall"), (r"∃", r"\\exists"),
            (r"→", r"\\to"), (r"⇒", r"\\implies"), (r"±", r"\\pm"),
        ]
        
        normalized = text
        for char, latex in replacements:
            normalized = normalized.replace(char, latex)
            
        return normalized

    def _detect_math_heavy_pages(self, doc, threshold: int = 3) -> List[int]:
        """
        Identify pages that contain significant math content.
        Returns a list of 1-based page numbers.
        """
        math_pages = set()
        math_symbols = set("∑∫√∂∇≈≤≥∞αβγθπµσΩ∈∀∃⇒±")
        
        # Efficient pass: Iterate all items once
        page_math_scores = {}  # page_no -> score
        
        for item, _ in doc.iterate_items():
            # Get page number (1-based)
            page_no = 1
            if hasattr(item, "prov") and item.prov:
                # prov is a list of Provenance items
                page_no = item.prov[0].page_no
            
            # Check for Formula label
            # Docling label enum or string: "formula", "equation"
            label = getattr(item, "label", "").lower() if hasattr(item, "label") else ""
            if "formula" in label or "equation" in label:
                page_math_scores[page_no] = page_math_scores.get(page_no, 0) + 10  # High score for explicit label
                
            # Check text content
            text = getattr(item, "text", "")
            if text:
                # Unicode density check
                symbol_count = sum(1 for char in text if char in math_symbols)
                
                # Superscript/Subscript check
                # ranges: super (²³¹⁰...): \u00B2, \u00B3, \u00B9, \u2070-\u207F
                # sub (₀₁...): \u2080-\u209C
                sub_super_count = 0
                for char in text:
                    code = ord(char)
                    if (0x2070 <= code <= 0x207F) or (0x2080 <= code <= 0x209C) or code in [0xB2, 0xB3, 0xB9]:
                        sub_super_count += 1
                
                page_math_scores[page_no] = page_math_scores.get(page_no, 0) + symbol_count + (sub_super_count * 0.5)

        # Filter pages exceeding threshold
        for page_no, score in page_math_scores.items():
            if score >= threshold:
                math_pages.add(page_no)
                
        return sorted(list(math_pages))

    def _is_enriched_page_valid(self, doc, page_no: int) -> bool:
        """
        Check if an enriched page has valid output (detect garbled text).
        """
        # Get text for specific page from the doc
        page_text = ""
        # iterate_items(page_no) is supported
        for item, _ in doc.iterate_items(page_no=page_no):
             page_text += getattr(item, "text", "") + " "
                 
        if not page_text.strip():
            return True # Empty page is "valid" in the sense of not garbled
            
        # Check for garble markers
        if "/C0" in page_text or "/C1" in page_text:
            return False
            
        return True
